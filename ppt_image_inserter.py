#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 이미지 자동 삽입 스크립트
89개 업체의 이미지를 PPT에 자동으로 삽입합니다.
"""

import os
import glob
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import openpyxl
import io

class PPTImageInserter:
    def __init__(self, template_ppt_path, base_image_dir, output_ppt_path, excel_path=None):
        """
        Args:
            template_ppt_path: 템플릿 PPT 파일 경로
            base_image_dir: 이미지가 있는 기본 디렉토리 (예: C:/Users/user/Downloads/서울/중구)
            output_ppt_path: 출력 PPT 파일 경로
            excel_path: 엑셀 파일 경로 (업체 순서 정보)
        """
        self.template_ppt_path = template_ppt_path
        self.base_image_dir = base_image_dir
        self.output_ppt_path = output_ppt_path
        self.excel_path = excel_path
        
    def load_shop_order_from_excel(self):
        """엑셀 파일에서 업체 순서를 읽어옵니다."""
        if not self.excel_path or not os.path.exists(self.excel_path):
            print("엑셀 파일을 찾을 수 없습니다. 기본 정렬을 사용합니다.")
            return None
        
        try:
            wb = openpyxl.load_workbook(self.excel_path)
            ws = wb.active
            
            shop_order = []
            for row in range(2, ws.max_row + 1):
                shop_name = ws.cell(row, 3).value  # 매장명은 3번째 컬럼
                if shop_name:
                    shop_order.append(shop_name.strip())
            
            print(f"엑셀에서 {len(shop_order)}개 업체 순서 로드 완료")
            return shop_order
            
        except Exception as e:
            print(f"엑셀 파일 읽기 오류: {e}")
            return None
    
    def find_shop_directories(self):
        """업체 디렉토리 목록을 찾고 엑셀 순서대로 정렬합니다."""
        shop_dirs_dict = {}
        
        # 기본 경로에서 하위 디렉토리 탐색 (서울/중구 등 지역별로 탐색)
        if os.path.exists(self.base_image_dir):
            # downloads 폴더 내의 지역 폴더들 탐색
            for region in os.listdir(self.base_image_dir):
                region_path = os.path.join(self.base_image_dir, region)
                if os.path.isdir(region_path):
                    # 지역 폴더 내의 상세지역 폴더들 탐색
                    for detail_region in os.listdir(region_path):
                        detail_region_path = os.path.join(region_path, detail_region)
                        if os.path.isdir(detail_region_path):
                            # 상세지역 폴더 내의 업체 폴더들 탐색
                            for shop_name in os.listdir(detail_region_path):
                                shop_path = os.path.join(detail_region_path, shop_name)
                                if os.path.isdir(shop_path):
                                    # 업체 폴더 내에 업체 하위 폴더가 있는지 확인
                                    company_folder = os.path.join(shop_path, "업체")
                                    if os.path.exists(company_folder):
                                        # 네이버플레이스 캡처 이미지
                                        naver_capture = os.path.join(company_folder, "네이버플레이스_캡처.png")
                                        # 가격표 이미지들
                                        price_images = sorted(glob.glob(os.path.join(company_folder, "가격표_*.jpg")))
                                        if not price_images:
                                            price_images = sorted(glob.glob(os.path.join(company_folder, "가격표_*.png")))
                                        # 업체 이미지들
                                        image_files = glob.glob(os.path.join(company_folder, "업체_*.jpg"))
                                        
                                        if os.path.exists(naver_capture) or price_images or image_files:
                                            shop_dirs_dict[shop_name] = {
                                                'name': shop_name,
                                                'path': shop_path,
                                                'naver_capture': naver_capture if os.path.exists(naver_capture) else None,
                                                'price_images': price_images,
                                                'images': sorted(image_files)
                                            }
        
        # 엑셀 파일에서 순서 로드
        shop_order = self.load_shop_order_from_excel()
        
        if shop_order:
            # 엑셀 순서대로 정렬
            ordered_shops = []
            for shop_name in shop_order:
                if shop_name in shop_dirs_dict:
                    ordered_shops.append(shop_dirs_dict[shop_name])
                else:
                    print(f"  경고: '{shop_name}' 폴더를 찾을 수 없습니다.")
            
            # 엑셀에 없는 업체는 마지막에 추가
            for shop_name, shop_info in sorted(shop_dirs_dict.items()):
                if shop_name not in shop_order:
                    print(f"  추가: '{shop_name}' (엑셀에 없는 업체)")
                    ordered_shops.append(shop_info)
            
            return ordered_shops
        else:
            # 엑셀이 없으면 알파벳 순서로 정렬
            return [shop_dirs_dict[key] for key in sorted(shop_dirs_dict.keys())]
    
    def get_image_dimensions(self, image_path):
        """이미지의 크기를 확인합니다."""
        try:
            with Image.open(image_path) as img:
                return img.size  # (width, height)
        except Exception as e:
            print(f"이미지 읽기 오류 ({image_path}): {e}")
            return None
    
    def crop_image_to_square(self, image_path):
        """이미지를 중앙 기준으로 정사각형으로 크롭합니다."""
        try:
            with Image.open(image_path) as img:
                width, height = img.size
                
                # 정사각형 크기 결정 (짧은 쪽 기준)
                size = min(width, height)
                
                # 중앙 크롭 좌표 계산
                left = (width - size) // 2
                top = (height - size) // 2
                right = left + size
                bottom = top + size
                
                # 크롭
                cropped = img.crop((left, top, right, bottom))
                
                # 메모리에 저장
                img_byte_arr = io.BytesIO()
                cropped.save(img_byte_arr, format='PNG')
                img_byte_arr.seek(0)
                
                return img_byte_arr
        except Exception as e:
            print(f"이미지 크롭 오류 ({image_path}): {e}")
            return None
    
    def add_price_images_to_slide(self, prs, slide, price_images):
        """가격표 이미지를 슬라이드에 배치합니다 (최대 3개)."""
        # 안전 여백 설정
        margin_left = Inches(0.5)
        margin_top = Inches(0.5)
        margin_bottom = Inches(0.5)
        
        # 사용 가능한 영역
        available_width = prs.slide_width - margin_left * 2
        available_height = prs.slide_height - margin_top - margin_bottom
        
        num_images = len(price_images)
        
        if num_images == 1:
            # 가격표 1개: 중앙에 크게 배치
            img_path = price_images[0]
            img_dims = self.get_image_dimensions(img_path)
            if img_dims:
                img_width, img_height = img_dims
                aspect_ratio = img_width / img_height
                
                # 9:16 세로 비율에 맞춤
                height = available_height
                width = height * aspect_ratio
                
                # 너무 넓으면 너비에 맞춤
                if width > available_width:
                    width = available_width
                    height = width / aspect_ratio
                
                # 중앙 정렬
                left = margin_left + (available_width - width) / 2
                top = margin_top + (available_height - height) / 2
                
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        
        else:
            # 가격표 2-3개: 가로로 나란히 배치
            gap = Inches(0.3)  # 이미지 간 간격
            
            # 각 이미지 너비 계산
            total_gap = gap * (num_images - 1)
            img_width = (available_width - total_gap) / num_images
            
            # 9:16 비율로 높이 계산
            img_height = img_width * (16 / 9)
            
            # 높이가 너무 크면 높이에 맞춤
            if img_height > available_height:
                img_height = available_height
                img_width = img_height * (9 / 16)
                # 너비 재계산
                total_gap = gap * (num_images - 1)
                total_width = img_width * num_images + total_gap
            else:
                total_width = available_width
            
            # 시작 위치 계산 (중앙 정렬)
            start_left = margin_left + (available_width - total_width) / 2
            top = margin_top + (available_height - img_height) / 2
            
            for i, img_path in enumerate(price_images):
                left = start_left + (i * (img_width + gap))
                slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)
    
    def add_images_to_slide(self, prs, slide, images_to_add, slide_height):
        """슬라이드에 이미지를 최적으로 배치합니다."""
        # 안전 여백 설정
        margin_left = Inches(0.5)
        margin_right = Inches(0.5)
        margin_top = Inches(1.2)
        margin_bottom = Inches(0.3)
        
        # 사용 가능한 영역
        available_width = prs.slide_width - margin_left - margin_right
        available_height = prs.slide_height - margin_top - margin_bottom
        
        if len(images_to_add) == 1:
            # 이미지 1개: 정사각형으로 크롭하여 중앙 배치
            img_path = images_to_add[0]
            
            # 사용 가능한 영역의 작은 쪽에 맞춤
            square_size = min(available_width, available_height)
            
            # 이미지를 정사각형으로 크롭
            cropped_img = self.crop_image_to_square(img_path)
            if cropped_img:
                # 중앙 정렬
                left = margin_left + (available_width - square_size) / 2
                top = margin_top + (available_height - square_size) / 2
                
                slide.shapes.add_picture(cropped_img, left, top, width=square_size, height=square_size)
        
        elif len(images_to_add) == 3:
            # 세로 이미지 3개: 가로로 나란히 배치 (안전 여백 포함)
            gap = Inches(0.3)  # 이미지 간 간격
            
            # 정사각형 크기 계산 (사용 가능한 너비에서 계산)
            square_size = (available_width - gap * 2) / 3
            
            # 높이 체크 - 너무 크면 높이에 맞춤
            if square_size > available_height:
                square_size = available_height
            
            # 시작 위치 계산 (중앙 정렬)
            total_width = square_size * 3 + gap * 2
            start_left = margin_left + (available_width - total_width) / 2
            top = margin_top + (available_height - square_size) / 2
            
            for i, img_path in enumerate(images_to_add):
                # 이미지를 정사각형으로 크롭
                cropped_img = self.crop_image_to_square(img_path)
                if cropped_img:
                    left = start_left + (i * (square_size + gap))
                    slide.shapes.add_picture(cropped_img, left, top, width=square_size, height=square_size)
    
    def add_shop_to_ppt(self, prs, shop_info, shop_number):
        """업체 정보를 PPT에 추가합니다."""
        naver_capture = shop_info.get('naver_capture')
        price_images = shop_info.get('price_images', [])
        images = shop_info['images']
        shop_name = shop_info['name']
        
        slide_layout = prs.slide_layouts[6]  # Blank layout (placeholder 없음)
        
        # 1. 표지 슬라이드 (업체명 + 네이버플레이스 캡처)
        slide = prs.slides.add_slide(slide_layout)
        
        # 제목 텍스트 박스 추가 (샘플 형식대로)
        textbox = slide.shapes.add_textbox(Inches(0), Inches(0.4), Inches(2.08), Inches(0.7))
        text_frame = textbox.text_frame
        text_frame.text = shop_name
        text_frame.word_wrap = True
        
        # 텍스트 스타일 설정: 맑은 고딕, 18pt
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = '맑은 고딕'
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
        
        # 네이버플레이스 캡처 이미지 추가 (원본 비율 유지, 자르지 않음)
        if naver_capture:
            try:
                img_dims = self.get_image_dimensions(naver_capture)
                if img_dims:
                    img_width, img_height = img_dims
                    aspect_ratio = img_width / img_height
                    
                    # 슬라이드에 맞게 크기 조정 (자르지 않고 전체 표시)
                    max_width = Inches(12)
                    max_height = Inches(6)
                    
                    if aspect_ratio > max_width / max_height:
                        width = max_width
                        height = width / aspect_ratio
                    else:
                        height = max_height
                        width = height * aspect_ratio
                    
                    left = (prs.slide_width - width) / 2
                    top = Inches(1.2)
                    
                    slide.shapes.add_picture(naver_capture, left, top, width=width, height=height)
                    print(f"  - 표지 슬라이드 추가: {shop_name} (네이버플레이스 캡처 포함)")
            except Exception as e:
                print(f"  - 네이버플레이스 캡처 추가 실패: {e}")
                print(f"  - 표지 슬라이드 추가: {shop_name} (텍스트만)")
        else:
            print(f"  - 표지 슬라이드 추가: {shop_name}")
        
        # 2. 가격표 슬라이드 추가 (가격표_*.jpg/png, 3개씩 배치)
        if price_images:
            idx = 0
            while idx < len(price_images):
                slide = prs.slides.add_slide(slide_layout)
                
                # 최대 3개씩 묶어서 배치
                batch_images = price_images[idx:idx+3]
                
                try:
                    self.add_price_images_to_slide(prs, slide, batch_images)
                    print(f"  - 가격표 슬라이드 추가: {len(batch_images)}개 이미지")
                except Exception as e:
                    print(f"  - 가격표 추가 실패: {e}")
                
                idx += 3
        
        # 3. 업체 이미지 슬라이드 추가 (업체_*.jpg) - 이미지를 그룹으로 묶어 처리
        idx = 0
        while idx < len(images):
            slide = prs.slides.add_slide(slide_layout)
            
            # 이미지 분석 및 배치
            try:
                # 현재 이미지와 다음 이미지들 확인
                current_img = images[idx]
                img_dims = self.get_image_dimensions(current_img)
                
                if img_dims:
                    img_width, img_height = img_dims
                    aspect_ratio = img_width / img_height
                    
                    # 가로 이미지면 1개만, 세로 이미지면 3개 배치
                    if aspect_ratio > 1.3:
                        # 가로 이미지: 1개만
                        images_to_add = [current_img]
                        idx += 1
                    else:
                        # 세로 이미지: 최대 3개 배치
                        images_to_add = [current_img]
                        
                        # 다음 이미지 확인 (최대 2개 더)
                        for j in range(1, 3):
                            if idx + j < len(images):
                                next_dims = self.get_image_dimensions(images[idx + j])
                                if next_dims and next_dims[0] / next_dims[1] <= 1.3:
                                    images_to_add.append(images[idx + j])
                                else:
                                    break
                            else:
                                break
                        
                        idx += len(images_to_add)
                    
                    self.add_images_to_slide(prs, slide, images_to_add, prs.slide_height)
                    print(f"  - 이미지 슬라이드 추가: {len(images_to_add)}개 이미지")
                else:
                    idx += 1
                    
            except Exception as e:
                print(f"  - 이미지 추가 실패: {e}")
                idx += 1
    
    def create_ppt(self, sample_mode=False, sample_count=10):
        """전체 PPT를 생성합니다.
        
        Args:
            sample_mode: 샘플 모드 여부 (True: 일부만 생성, False: 전체 생성)
            sample_count: 샘플 모드일 때 생성할 업체 수
        """
        print("=" * 60)
        if sample_mode:
            print(f"PPT 샘플 생성 시작 (최대 {sample_count}개 업체)")
        else:
            print("PPT 전체 생성 시작")
        print("=" * 60)
        
        # 업체 디렉토리 찾기
        print(f"\n이미지 디렉토리 스캔 중: {self.base_image_dir}")
        shop_dirs = self.find_shop_directories()
        
        if not shop_dirs:
            print(f"오류: 업체 디렉토리를 찾을 수 없습니다.")
            print(f"경로를 확인하세요: {self.base_image_dir}")
            return False
        
        print(f"발견된 업체 수: {len(shop_dirs)}")
        
        # 샘플 모드일 경우 업체 수 제한
        if sample_mode:
            original_count = len(shop_dirs)
            shop_dirs = shop_dirs[:sample_count]
            print(f"샘플 모드: {original_count}개 중 {len(shop_dirs)}개 업체만 처리합니다.")
        
        # 템플릿 PPT 로드
        print(f"\n템플릿 PPT 로드 중: {self.template_ppt_path}")
        prs = Presentation(self.template_ppt_path)
        
        # 표지 슬라이드 추가
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        textbox = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(5.33), Inches(1.5))
        text_frame = textbox.text_frame
        if sample_mode:
            text_frame.text = f"세신샵 업체 정보 (샘플 {len(shop_dirs)}개)"
        else:
            text_frame.text = "세신샵 업체 정보"
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Inches(0.6)
            paragraph.font.bold = True
        
        print("\n업체별 슬라이드 생성 중...")
        
        # 각 업체별로 슬라이드 추가
        for idx, shop_info in enumerate(shop_dirs, 1):
            print(f"\n[{idx}/{len(shop_dirs)}] {shop_info['name']} 처리 중...")
            print(f"  이미지 수: {len(shop_info['images'])}")
            self.add_shop_to_ppt(prs, shop_info, idx)
        
        # PPT 저장
        print(f"\nPPT 저장 중: {self.output_ppt_path}")
        prs.save(self.output_ppt_path)
        
        print("=" * 60)
        if sample_mode:
            print(f"샘플 완료! {len(shop_dirs)}개 업체의 슬라이드가 생성되었습니다.")
        else:
            print(f"전체 완료! 총 {len(shop_dirs)}개 업체의 슬라이드가 생성되었습니다.")
        print(f"출력 파일: {self.output_ppt_path}")
        print("=" * 60)
        
        return True


def main():
    """메인 실행 함수"""
    
    # ========================================
    # 설정 값 (필요에 따라 수정하세요)
    # ========================================
    
    # 스크립트가 있는 현재 디렉토리 경로
    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
    
    # 이미지 디렉토리 (naver-map-photo-downloader의 downloads 폴더)
    BASE_IMAGE_DIR = r"C:\Users\user\Documents\GitHub\naver-map-photo-downloader\downloads"
    
    # 템플릿 PPT 파일 (스크립트와 같은 폴더)
    TEMPLATE_PPT = os.path.join(SCRIPT_DIR, "세신샵.pptx")
    
    # 엑셀 파일 (업체 순서 정보, 스크립트와 같은 폴더)
    EXCEL_FILE = os.path.join(SCRIPT_DIR, "리스트_네이버지도링크추가 - 복사본.xlsx")
    
    # 출력 PPT 파일 (스크립트와 같은 폴더)
    OUTPUT_PPT = os.path.join(SCRIPT_DIR, "세신샵_완성본.pptx")
    
    # ========================================
    
    # 경로 확인
    if not os.path.exists(TEMPLATE_PPT):
        print(f"오류: 템플릿 PPT 파일을 찾을 수 없습니다: {TEMPLATE_PPT}")
        return
    
    # 작업 모드 선택
    print("\n" + "=" * 60)
    print("PPT 자동 생성 프로그램")
    print("=" * 60)
    print("\n작업 모드를 선택하세요:")
    print("  1. 샘플 모드 (처음 10개 업체만 생성)")
    print("  2. 전체 작업 (모든 업체 생성)")
    print("=" * 60)
    
    while True:
        try:
            choice = input("\n선택 (1 또는 2): ").strip()
            
            if choice == "1":
                # 샘플 모드
                sample_mode = True
                sample_count = 10
                output_file = OUTPUT_PPT.replace(".pptx", "_샘플.pptx")
                print(f"\n✓ 샘플 모드 선택됨 (처음 {sample_count}개 업체)")
                break
            elif choice == "2":
                # 전체 작업
                sample_mode = False
                sample_count = 0
                output_file = OUTPUT_PPT
                print("\n✓ 전체 작업 모드 선택됨 (모든 업체)")
                break
            else:
                print("❌ 잘못된 입력입니다. 1 또는 2를 입력하세요.")
        except KeyboardInterrupt:
            print("\n\n프로그램을 종료합니다.")
            return
    
    # 확인 메시지
    print("\n" + "=" * 60)
    print("설정 확인:")
    print(f"  - 모드: {'샘플 (10개)' if sample_mode else '전체 작업'}")
    print(f"  - 이미지 디렉토리: {BASE_IMAGE_DIR}")
    print(f"  - 엑셀 파일: {EXCEL_FILE}")
    print(f"  - 출력 파일: {output_file}")
    print("=" * 60)
    
    proceed = input("\n진행하시겠습니까? (y/n): ").strip().lower()
    if proceed != 'y':
        print("작업을 취소합니다.")
        return
    
    # PPT 생성기 실행
    inserter = PPTImageInserter(
        template_ppt_path=TEMPLATE_PPT,
        base_image_dir=BASE_IMAGE_DIR,
        output_ppt_path=output_file,
        excel_path=EXCEL_FILE
    )
    
    inserter.create_ppt(sample_mode=sample_mode, sample_count=sample_count)


if __name__ == "__main__":
    main()
